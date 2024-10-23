# app.py

import streamlit as st
import os
import uuid
import io
from dotenv import load_dotenv
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance,
    VectorParams,
    PointStruct,
)
from pdfminer.high_level import extract_text_to_fp
from pdfminer.layout import LAParams
from docx import Document
from pptx import Presentation  # For PPTX support
import pandas as pd  # For Excel support
from sentence_transformers import SentenceTransformer
import numpy as np
import tempfile
from typing import List, Dict, Tuple
import logging
from pdf2image import convert_from_bytes
from PIL import Image
import base64
import re
from transformers import pipeline, AutoModelForSeq2SeqLM, AutoTokenizer

# =========================
# Logging Configuration
# =========================

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler("app.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# =========================
# Configuration and Setup
# =========================

# Load environment variables from .env file
load_dotenv()

# =========================
# Streamlit Page Configuration
# =========================

# Set Streamlit page configuration as the very first Streamlit command
st.set_page_config(
    page_title="📄 RAG App with Qdrant Cloud",
    layout="wide",
    initial_sidebar_state="expanded",
)

# =========================
# Retrieve Environment Variables
# =========================

QDRANT_API_KEY = os.getenv("QDRANT_API_KEY")
QDRANT_URL = os.getenv("QDRANT_URL")  # e.g., "https://your-instance.qdrant.io"

# =========================
# Validate Environment Variables
# =========================

if not all([QDRANT_API_KEY, QDRANT_URL]):
    st.error("Please ensure that QDRANT_API_KEY and QDRANT_URL are set in the .env file.")
    logger.error("Missing QDRANT_API_KEY or QDRANT_URL in .env file.")
    st.stop()

# =========================
# Initialize Qdrant Client
# =========================

try:
    qdrant_client = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY)
    logger.info("Successfully connected to Qdrant Cloud.")
except Exception as e:
    st.error(f"Failed to connect to Qdrant Cloud: {e}")
    logger.error(f"Failed to connect to Qdrant Cloud: {e}")
    st.stop()

# Define Qdrant collection name
COLLECTION_NAME = "documents"

# Define embedding model parameters
EMBEDDING_MODEL_NAME = "all-MiniLM-L6-v2"  # A lightweight, efficient model
VECTOR_SIZE = 384  # Dimension size for all-MiniLM-L6-v2

# =========================
# Initialize Sentence Transformer Model
# =========================

@st.cache_resource
def load_embedding_model():
    """
    Loads the pre-trained SentenceTransformer model for generating embeddings.

    Returns:
        SentenceTransformer: The loaded embedding model.

    Example:
        model = load_embedding_model()
    """
    try:
        model = SentenceTransformer(EMBEDDING_MODEL_NAME)
        logger.info(f"Loaded embedding model: {EMBEDDING_MODEL_NAME}")
        return model
    except Exception as e:
        st.error(f"Failed to load embedding model: {e}")
        logger.error(f"Failed to load embedding model: {e}")
        st.stop()

embedding_model = load_embedding_model()

# =========================
# Initialize Open-Source LLM
# =========================

@st.cache_resource
def load_llm_pipeline():
    """
    Loads the pre-trained open-source LLM pipeline for answer generation.

    Returns:
        transformers.Pipeline: The loaded LLM pipeline.
    """
    try:
        model_name = "google/flan-t5-base"  # You can choose other models like "EleutherAI/gpt-j-6B"
        tokenizer = AutoTokenizer.from_pretrained(model_name)
        model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
        llm_pipeline = pipeline("text2text-generation", model=model, tokenizer=tokenizer)
        logger.info(f"Loaded LLM model: {model_name}")
        return llm_pipeline
    except Exception as e:
        st.error(f"Failed to load LLM model: {e}")
        logger.error(f"Failed to load LLM model: {e}")
        st.stop()

llm_pipeline = load_llm_pipeline()

# =========================
# Utility Functions
# =========================

def initialize_qdrant_collection():
    """
    Initializes the Qdrant collection with the specified schema.
    If the collection already exists, it will not be recreated.

    Example:
        initialize_qdrant_collection()

    Output:
        None
    """
    try:
        existing_collections = qdrant_client.get_collections().collections
        if COLLECTION_NAME not in [col.name for col in existing_collections]:
            qdrant_client.create_collection(
                collection_name=COLLECTION_NAME,
                vectors_config=VectorParams(
                    size=VECTOR_SIZE,
                    distance=Distance.COSINE,
                ),
            )
            st.sidebar.success(f"Qdrant collection '{COLLECTION_NAME}' created successfully.")
            logger.info(f"Qdrant collection '{COLLECTION_NAME}' created successfully.")
        else:
            st.sidebar.info(f"Qdrant collection '{COLLECTION_NAME}' already exists.")
            logger.info(f"Qdrant collection '{COLLECTION_NAME}' already exists.")
    except Exception as e:
        st.sidebar.error(f"Error initializing Qdrant collection: {e}")
        logger.error(f"Error initializing Qdrant collection: {e}")
        st.stop()

def extract_text_from_pdf(file: io.BytesIO) -> List[Tuple[int, str]]:
    """
    Extracts text from a PDF file on a per-page basis.

    Args:
        file (io.BytesIO): The uploaded PDF file.

    Returns:
        List[Tuple[int, str]]: A list of tuples containing page numbers and extracted text.

    Example:
        pages = extract_text_from_pdf(uploaded_pdf)
    """
    try:
        pages = []
        laparams = LAParams()
        with tempfile.TemporaryFile() as temp_file:
            extract_text_to_fp(file, temp_file, laparams=laparams, output_type='text', codec=None)
            temp_file.seek(0)
            text = temp_file.read().decode('utf-8')
            # Split text by form feed character '\f' to separate pages
            page_texts = text.split('\f')
            for i, page in enumerate(page_texts, start=1):
                clean_page = page.strip()
                if clean_page:
                    pages.append((i, clean_page))
        logger.info(f"Extracted text from PDF with {len(pages)} pages.")
        return pages
    except Exception as e:
        st.warning(f"Failed to extract text from PDF: {e}")
        logger.error(f"Failed to extract text from PDF: {e}")
        return []

def extract_text_from_docx(file: io.BytesIO) -> List[Tuple[int, str]]:
    """
    Extracts text from a DOCX file on a per-paragraph basis.

    Args:
        file (io.BytesIO): The uploaded DOCX file.

    Returns:
        List[Tuple[int, str]]: A list of tuples containing paragraph numbers and extracted text.

    Example:
        paragraphs = extract_text_from_docx(uploaded_docx)
    """
    try:
        doc = Document(io.BytesIO(file.read()))
        paragraphs = []
        for i, para in enumerate(doc.paragraphs, start=1):
            clean_para = para.text.strip()
            if clean_para:
                paragraphs.append((i, clean_para))
        logger.info(f"Extracted text from DOCX with {len(paragraphs)} paragraphs.")
        return paragraphs
    except Exception as e:
        st.warning(f"Failed to extract text from DOCX: {e}")
        logger.error(f"Failed to extract text from DOCX: {e}")
        return []

def extract_text_from_excel(file: io.BytesIO) -> List[Tuple[int, str]]:
    """
    Extracts text from an Excel file on a per-cell basis.

    Args:
        file (io.BytesIO): The uploaded Excel file.

    Returns:
        List[Tuple[int, str]]: A list of tuples containing cell indices and extracted text.

    Example:
        cells = extract_text_from_excel(uploaded_excel)
    """
    try:
        dfs = pd.read_excel(file, sheet_name=None)
        cells = []
        for sheet_name, df in dfs.items():
            for row_idx, row in df.iterrows():
                for col_idx, cell_value in row.items():
                    if pd.notnull(cell_value):
                        clean_cell = str(cell_value).strip()
                        if clean_cell:
                            cell_identifier = f"{sheet_name}!{col_idx}{row_idx + 2}"  # Excel rows are 1-indexed
                            cells.append((cell_identifier, clean_cell))
        logger.info(f"Extracted text from Excel with {len(cells)} cells.")
        return cells
    except Exception as e:
        st.warning(f"Failed to extract text from Excel: {e}")
        logger.error(f"Failed to extract text from Excel: {e}")
        return []

def extract_text_from_pptx(file: io.BytesIO) -> List[Tuple[int, str]]:
    """
    Extracts text from a PowerPoint file on a per-slide basis.

    Args:
        file (io.BytesIO): The uploaded PowerPoint file.

    Returns:
        List[Tuple[int, str]]: A list of tuples containing slide numbers and extracted text.

    Example:
        slides = extract_text_from_pptx(uploaded_pptx)
    """
    try:
        prs = Presentation(file)
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + " "
            clean_slide = slide_text.strip()
            if clean_slide:
                slides.append((i, clean_slide))
        logger.info(f"Extracted text from PowerPoint with {len(slides)} slides.")
        return slides
    except Exception as e:
        st.warning(f"Failed to extract text from PowerPoint: {e}")
        logger.error(f"Failed to extract text from PowerPoint: {e}")
        return []

def extract_text_from_txt(file: io.BytesIO) -> List[Tuple[int, str]]:
    """
    Extracts text from a text file.

    Args:
        file (io.BytesIO): The uploaded text file.

    Returns:
        List[Tuple[int, str]]: A list containing a single tuple with the entire text.

    Example:
        text = extract_text_from_txt(uploaded_txt)
    """
    try:
        text = file.read().decode('utf-8')
        lines = text.split('\n')
        paragraphs = []
        for i, line in enumerate(lines, start=1):
            clean_line = line.strip()
            if clean_line:
                paragraphs.append((i, clean_line))
        logger.info(f"Extracted text from text file with {len(paragraphs)} lines.")
        return paragraphs
    except Exception as e:
        st.warning(f"Failed to extract text from text file: {e}")
        logger.error(f"Failed to extract text from text file: {e}")
        return []

def extract_text_from_file(file) -> List[Dict]:
    """
    Determines the file type and extracts text accordingly, returning a list of chunks with metadata.

    Args:
        file (UploadedFile): The uploaded file from Streamlit.

    Returns:
        List[Dict]: A list of dictionaries containing chunk text and metadata.

    Example:
        chunks = extract_text_from_file(uploaded_file)
    """
    chunks = []
    if file.type == "application/pdf":
        pages = extract_text_from_pdf(file)
        for page_num, text in pages:
            page_chunks = chunk_text(text, max_length=500)
            for chunk in page_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "pdf",
                    "page_number": page_num,
                    "chunk_text": chunk
                })
    elif file.type in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ]:
        paragraphs = extract_text_from_docx(file)
        for para_num, text in paragraphs:
            para_chunks = chunk_text(text, max_length=500)
            for chunk in para_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "docx",
                    "paragraph_number": para_num,
                    "chunk_text": chunk
                })
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        slides = extract_text_from_pptx(file)
        for slide_num, text in slides:
            slide_chunks = chunk_text(text, max_length=500)
            for chunk in slide_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "pptx",
                    "slide_number": slide_num,
                    "chunk_text": chunk
                })
    elif file.type in [
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ]:
        cells = extract_text_from_excel(file)
        for cell_id, text in cells:
            cell_chunks = chunk_text(text, max_length=500)
            for chunk in cell_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "xlsx",
                    "cell_id": cell_id,
                    "chunk_text": chunk
                })
    elif file.type == "text/plain":
        lines = extract_text_from_txt(file)
        for line_num, text in lines:
            line_chunks = chunk_text(text, max_length=500)
            for chunk in line_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "txt",
                    "line_number": line_num,
                    "chunk_text": chunk
                })
    else:
        st.warning(f"Unsupported file type: {file.type}")
        logger.warning(f"Unsupported file type: {file.type}")
    logger.info(f"Extracted {len(chunks)} chunks from '{file.name}'.")
    return chunks

def chunk_text(text: str, max_length: int = 500) -> List[str]:
    """
    Splits text into chunks based on sentence delimiters, ensuring each chunk does not exceed max_length.

    Args:
        text (str): The input text to split.
        max_length (int): Maximum number of characters per chunk.

    Returns:
        List[str]: A list of text chunks.

    Example:
        chunks = chunk_text("Your long text here.", max_length=500)
    """
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = ""
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        if len(current_chunk) + len(sentence) + 1 <= max_length:
            current_chunk += f" {sentence}"
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = f"{sentence}"
    if current_chunk:
        chunks.append(current_chunk.strip())
    logger.debug(f"Chunked text into {len(chunks)} chunks.")
    return chunks

def generate_embedding(text: str) -> list:
    """
    Generates an embedding vector for the given text using SentenceTransformer.

    Args:
        text (str): The input text to generate embedding for.

    Returns:
        list: The embedding vector.

    Example:
        embedding = generate_embedding("Sample text for embedding.")
    """
    try:
        embedding = embedding_model.encode(text)
        # Normalize the embedding to unit vector for cosine similarity
        embedding = embedding / np.linalg.norm(embedding)
        logger.debug("Generated and normalized embedding.")
        return embedding.tolist()
    except Exception as e:
        st.warning(f"Failed to generate embedding: {e}")
        logger.error(f"Failed to generate embedding: {e}")
        return None

def store_document(chunk: Dict):
    """
    Stores a chunk's embedding and metadata in Qdrant.

    Args:
        chunk (Dict): A dictionary containing chunk text and metadata.

    Returns:
        None

    Example:
        store_document({"file_name": "doc.pdf", "document_id": "1234", "page_number": 1, "chunk_text": "Sample text."})
    """
    embedding = generate_embedding(chunk["chunk_text"])
    if embedding is None:
        st.sidebar.error(f"Failed to generate embedding for chunk in '{chunk['file_name']}'.")
        logger.error(f"Failed to generate embedding for chunk in '{chunk['file_name']}'.")
        return
    payload = {
        "file_name": chunk["file_name"],
        "document_id": chunk["document_id"],
        "file_type": chunk["file_type"],
    }
    if chunk["file_type"] == "pdf":
        payload["page_number"] = chunk.get("page_number", "N/A")
    elif chunk["file_type"] == "docx":
        payload["paragraph_number"] = chunk.get("paragraph_number", "N/A")
    elif chunk["file_type"] == "pptx":
        payload["slide_number"] = chunk.get("slide_number", "N/A")
    elif chunk["file_type"] == "xlsx":
        payload["cell_id"] = chunk.get("cell_id", "N/A")
    elif chunk["file_type"] == "txt":
        payload["line_number"] = chunk.get("line_number", "N/A")
    payload["chunk_text"] = chunk["chunk_text"]

    # Use a single UUID for point ID
    point_id = str(uuid.uuid4())

    point = PointStruct(
        id=point_id,
        vector=embedding,
        payload=payload
    )
    try:
        qdrant_client.upsert(collection_name=COLLECTION_NAME, points=[point])
        st.sidebar.success(f"Indexed a chunk from '{chunk['file_name']}'.")
        logger.info(f"Indexed chunk ID: {point_id} from '{chunk['file_name']}'.")
    except Exception as e:
        st.sidebar.error(f"Failed to store chunk in Qdrant: {e}")
        logger.error(f"Failed to store chunk ID: {point_id} in Qdrant: {e}")

def query_similar_documents(query_embedding: list, top_k: int = 5):
    """
    Queries Qdrant for documents similar to the provided embedding.

    Args:
        query_embedding (list): The embedding vector of the query.
        top_k (int): Number of top similar documents to retrieve.

    Returns:
        list: A list of similar documents with their metadata and scores.

    Example:
        results = query_similar_documents([0.1, 0.2, ...], top_k=3)
    """
    try:
        search_result = qdrant_client.search(
            collection_name=COLLECTION_NAME,
            query_vector=query_embedding,
            limit=top_k,
        )
        logger.info(f"Queried Qdrant for top {top_k} similar documents.")
        return search_result
    except Exception as e:
        st.error(f"Failed to query Qdrant: {e}")
        logger.error(f"Failed to query Qdrant: {e}")
        return []

def get_file_icon(file_type: str) -> str:
    """
    Returns an emoji icon based on the file type.

    Args:
        file_type (str): The type of the file (e.g., 'pdf', 'docx').

    Returns:
        str: An emoji representing the file type.

    Example:
        icon = get_file_icon("pdf")  # Returns 📄
    """
    icons = {
        "pdf": "📄",
        "docx": "📃",
        "pptx": "📊",
        "xlsx": "📈",
        "txt": "📝",
        "default": "📄"
    }
    return icons.get(file_type, icons["default"])

def generate_pdf_thumbnail(file: io.BytesIO) -> bytes:
    """
    Generates a thumbnail image for the first page of a PDF.

    Args:
        file (io.BytesIO): The uploaded PDF file.

    Returns:
        bytes: The thumbnail image in bytes.

    Example:
        thumbnail = generate_pdf_thumbnail(uploaded_pdf)
    """
    try:
        # Reset file pointer
        file.seek(0)
        # Convert first page of PDF to image
        images = convert_from_bytes(file.read(), first_page=1, last_page=1, size=(200, 200))
        if images:
            img_buffer = io.BytesIO()
            images[0].save(img_buffer, format='PNG')
            img_bytes = img_buffer.getvalue()
            logger.info("Generated PDF thumbnail.")
            return img_bytes
        else:
            logger.warning("No images generated for PDF thumbnail.")
            return None
    except Exception as e:
        st.warning(f"Failed to generate thumbnail for PDF: {e}")
        logger.error(f"Failed to generate thumbnail for PDF: {e}")
        return None

def get_thumbnail(file_type: str, file: io.BytesIO) -> bytes:
    """
    Retrieves a thumbnail based on the file type.

    Args:
        file_type (str): The type of the file (e.g., 'pdf', 'docx').
        file (io.BytesIO): The uploaded file.

    Returns:
        bytes: The thumbnail image in bytes, or None.
    """
    if file_type == "pdf":
        return generate_pdf_thumbnail(file)
    elif file_type == "pptx":
        # For PPTX, you can implement similar logic or return a default icon
        return None
    else:
        return None

def generate_answer(query: str, context_chunks: List[str]) -> str:
    """
    Generates an answer to the user's query based on the context chunks using an open-source LLM.

    Args:
        query (str): The user's question.
        context_chunks (List[str]): List of text chunks from retrieved documents.

    Returns:
        str: The generated answer.

    Example:
        answer = generate_answer("What is AI?", ["AI stands for...", "It is used in..."])
    """
    try:
        # Combine context chunks into a single context string
        context = "\n\n".join(context_chunks)
        # Structure the prompt as per the user's requirement
        prompt = f"The info found in documents collected explain that:\n\n{context}\n\nQuestion: {query}\n\nAnswer:"
        # Use the open-source LLM pipeline to generate an answer
        response = llm_pipeline(prompt, max_length=200, num_return_sequences=1)
        answer = response[0]['generated_text'].strip()
        logger.info("Generated answer using open-source LLM.")
        return answer
    except Exception as e:
        st.error(f"Failed to generate answer: {e}")
        logger.error(f"Failed to generate answer: {e}")
        return "I'm sorry, I couldn't generate an answer at this time."

# =========================
# Streamlit Application
# =========================

def main():
    """
    The main function that defines the Streamlit app layout and interactions.

    Args:
        None

    Returns:
        None

    Example:
        main()
    """
    # Initialize Qdrant collection
    initialize_qdrant_collection()

    # App Title
    st.title("📄 Retrieval-Augmented Generation (RAG) App with Qdrant Cloud")

    # =========================
    # Sidebar: File Upload
    # =========================
    st.sidebar.header("📂 Upload Documents")
    uploaded_files = st.sidebar.file_uploader(
        "Choose files", type=["pdf", "docx", "xlsx", "pptx", "txt"], accept_multiple_files=True
    )

    if uploaded_files:
        for uploaded_file in uploaded_files:
            st.sidebar.write(f"Processing: {uploaded_file.name}")
            # Extract chunks from the uploaded file
            chunks = extract_text_from_file(uploaded_file)
            if not chunks:
                st.sidebar.warning(f"No text extracted from {uploaded_file.name}. Skipping.")
                continue

            # Store each chunk in Qdrant
            for chunk in chunks:
                store_document(chunk)

    st.sidebar.markdown("---")
    st.sidebar.info("Upload PDF, DOCX, XLSX, PPTX, or TXT files to index them for searching.")

    # =========================
    # Main Section: Querying
    # =========================
    st.header("🔍 Search Documents")

    # Input field for user query
    user_query = st.text_input("Enter your query here:", "")

    # Search button
    if st.button("Search"):
        if not user_query.strip():
            st.warning("Please enter a valid query.")
            logger.warning("User submitted an empty query.")
        else:
            with st.spinner("Generating embedding and searching for similar documents..."):
                # Generate embedding for the query
                query_embedding = generate_embedding(user_query)
                if query_embedding is None:
                    st.error("Failed to generate embedding for the query.")
                else:
                    # Query Qdrant for similar documents
                    results = query_similar_documents(query_embedding, top_k=5)

                    if not results:
                        st.info("No similar documents found.")
                    else:
                        st.success(f"Found {len(results)} similar document(s):")
                        context_chunks = []
                        for idx, result in enumerate(results, start=1):
                            payload = result.payload
                            file_name = payload.get("file_name", "Unknown")
                            document_id = payload.get("document_id", "N/A")
                            file_type = payload.get("file_type", "default")
                            chunk_text = payload.get("chunk_text", "")
                            page_number = payload.get("page_number", None)
                            paragraph_number = payload.get("paragraph_number", None)
                            slide_number = payload.get("slide_number", None)
                            cell_id = payload.get("cell_id", None)
                            line_number = payload.get("line_number", None)
                            score = result.score
                            thumbnail = payload.get("thumbnail", None)

                            icon = get_file_icon(file_type)

                            # Collect context chunks for answer generation
                            context_chunks.append(chunk_text)

                            # Display result with icon and metadata
                            with st.container():
                                col1, col2 = st.columns([1, 4])
                                with col1:
                                    if thumbnail:
                                        # Display thumbnail
                                        st.image(thumbnail, width=50)
                                    else:
                                        # Display icon
                                        st.markdown(icon)
                                with col2:
                                    st.markdown(f"**{file_name}**")
                                    st.markdown(f"**Document ID**: {document_id}")
                                    if file_type == "pdf" and page_number:
                                        st.markdown(f"**Page Number**: {page_number}")
                                    if file_type == "docx" and paragraph_number:
                                        st.markdown(f"**Paragraph Number**: {paragraph_number}")
                                    if file_type == "pptx" and slide_number:
                                        st.markdown(f"**Slide Number**: {slide_number}")
                                    if file_type == "xlsx" and cell_id:
                                        st.markdown(f"**Cell ID**: {cell_id}")
                                    if file_type == "txt" and line_number:
                                        st.markdown(f"**Line Number**: {line_number}")
                                    st.markdown(f"**Similarity Score**: {score:.4f}")
                                    st.markdown(f"**Snippet**: {chunk_text[:200]}...")  # Display first 200 chars
                                st.markdown("---")

                        # Generate an answer based on the retrieved context
                        with st.spinner("Generating answer..."):
                            answer = generate_answer(user_query, context_chunks)
                            st.header("📝 Answer")
                            st.write(answer)

    st.markdown("---")
    st.markdown("Developed with ❤️ using Streamlit and Qdrant Cloud.")

if __name__ == "__main__":
    main()
