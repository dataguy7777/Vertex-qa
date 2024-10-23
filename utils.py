# utils.py

import io
import tempfile
from typing import List, Tuple, Dict
from pdfminer.high_level import extract_text_to_fp
from pdfminer.layout import LAParams
from docx import Document
from pptx import Presentation
import pandas as pd
import re
from pdf2image import convert_from_bytes
from PIL import Image
from logger import logger
import streamlit as st

def extract_text_from_pdf(file: io.BytesIO) -> List[Tuple[int, str]]:
    """
    Extracts text from a PDF file on a per-page basis.
    """
    try:
        pages = []
        laparams = LAParams()
        with tempfile.TemporaryFile() as temp_file:
            extract_text_to_fp(file, temp_file, laparams=laparams, output_type='text', codec=None)
            temp_file.seek(0)
            text = temp_file.read().decode('utf-8')
            # Split text by form feed character '\f' to separate pages
            page_texts = text.split('\f')
            for i, page in enumerate(page_texts, start=1):
                clean_page = page.strip()
                if clean_page:
                    pages.append((i, clean_page))
        logger.info(f"Extracted text from PDF with {len(pages)} pages.")
        return pages
    except Exception as e:
        st.warning(f"Failed to extract text from PDF: {e}")
        logger.error(f"Failed to extract text from PDF: {e}")
        return []

def extract_text_from_docx(file: io.BytesIO) -> List[Tuple[int, str]]:
    """
    Extracts text from a DOCX file on a per-paragraph basis.
    """
    try:
        doc = Document(io.BytesIO(file.read()))
        paragraphs = []
        for i, para in enumerate(doc.paragraphs, start=1):
            clean_para = para.text.strip()
            if clean_para:
                paragraphs.append((i, clean_para))
        logger.info(f"Extracted text from DOCX with {len(paragraphs)} paragraphs.")
        return paragraphs
    except Exception as e:
        st.warning(f"Failed to extract text from DOCX: {e}")
        logger.error(f"Failed to extract text from DOCX: {e}")
        return []

def extract_text_from_excel(file: io.BytesIO) -> List[Tuple[str, str]]:
    """
    Extracts text from an Excel file on a per-cell basis.
    """
    try:
        dfs = pd.read_excel(file, sheet_name=None)
        cells = []
        for sheet_name, df in dfs.items():
            for row_idx, row in df.iterrows():
                for col_idx, cell_value in row.items():
                    if pd.notnull(cell_value):
                        clean_cell = str(cell_value).strip()
                        if clean_cell:
                            cell_identifier = f"{sheet_name}!{col_idx}{row_idx + 2}"  # Excel rows are 1-indexed
                            cells.append((cell_identifier, clean_cell))
        logger.info(f"Extracted text from Excel with {len(cells)} cells.")
        return cells
    except Exception as e:
        st.warning(f"Failed to extract text from Excel: {e}")
        logger.error(f"Failed to extract text from Excel: {e}")
        return []

def extract_text_from_pptx(file: io.BytesIO) -> List[Tuple[int, str]]:
    """
    Extracts text from a PowerPoint file on a per-slide basis.
    """
    try:
        prs = Presentation(file)
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + " "
            clean_slide = slide_text.strip()
            if clean_slide:
                slides.append((i, clean_slide))
        logger.info(f"Extracted text from PowerPoint with {len(slides)} slides.")
        return slides
    except Exception as e:
        st.warning(f"Failed to extract text from PowerPoint: {e}")
        logger.error(f"Failed to extract text from PowerPoint: {e}")
        return []

def extract_text_from_txt(file: io.BytesIO) -> List[Tuple[int, str]]:
    """
    Extracts text from a text file.
    """
    try:
        text = file.read().decode('utf-8')
        lines = text.split('\n')
        paragraphs = []
        for i, line in enumerate(lines, start=1):
            clean_line = line.strip()
            if clean_line:
                paragraphs.append((i, clean_line))
        logger.info(f"Extracted text from text file with {len(paragraphs)} lines.")
        return paragraphs
    except Exception as e:
        st.warning(f"Failed to extract text from text file: {e}")
        logger.error(f"Failed to extract text from text file: {e}")
        return []

def extract_text_from_file(file) -> List[Dict]:
    """
    Determines the file type and extracts text accordingly, returning a list of chunks with metadata.
    """
    chunks = []
    file_type = file.type

    if file_type == "application/pdf":
        pages = extract_text_from_pdf(file)
        for page_num, text in pages:
            page_chunks = chunk_text(text, max_length=500)
            for chunk in page_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "pdf",
                    "page_number": page_num,
                    "chunk_text": chunk
                })
    elif file_type in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ]:
        paragraphs = extract_text_from_docx(file)
        for para_num, text in paragraphs:
            para_chunks = chunk_text(text, max_length=500)
            for chunk in para_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "docx",
                    "paragraph_number": para_num,
                    "chunk_text": chunk
                })
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        slides = extract_text_from_pptx(file)
        for slide_num, text in slides:
            slide_chunks = chunk_text(text, max_length=500)
            for chunk in slide_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "pptx",
                    "slide_number": slide_num,
                    "chunk_text": chunk
                })
    elif file_type in [
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ]:
        cells = extract_text_from_excel(file)
        for cell_id, text in cells:
            cell_chunks = chunk_text(text, max_length=500)
            for chunk in cell_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "xlsx",
                    "cell_id": cell_id,
                    "chunk_text": chunk
                })
    elif file_type == "text/plain":
        lines = extract_text_from_txt(file)
        for line_num, text in lines:
            line_chunks = chunk_text(text, max_length=500)
            for chunk in line_chunks:
                chunks.append({
                    "file_name": file.name,
                    "document_id": str(uuid.uuid4()),
                    "file_type": "txt",
                    "line_number": line_num,
                    "chunk_text": chunk
                })
    else:
        st.warning(f"Unsupported file type: {file.type}")
        logger.warning(f"Unsupported file type: {file.type}")

    logger.info(f"Extracted {len(chunks)} chunks from '{file.name}'.")
    return chunks

def chunk_text(text: str, max_length: int = 500) -> List[str]:
    """
    Splits text into chunks based on sentence delimiters, ensuring each chunk does not exceed max_length.
    """
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = ""
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        if len(current_chunk) + len(sentence) + 1 <= max_length:
            current_chunk += f" {sentence}"
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = f"{sentence}"
    if current_chunk:
        chunks.append(current_chunk.strip())
    logger.debug(f"Chunked text into {len(chunks)} chunks.")
    return chunks

def get_file_icon(file_type: str) -> str:
    """
    Returns an emoji icon based on the file type.
    """
    icons = {
        "pdf": "📄",
        "docx": "📃",
        "pptx": "📊",
        "xlsx": "📈",
        "txt": "📝",
        "default": "📄"
    }
    return icons.get(file_type, icons["default"])

def generate_pdf_thumbnail(file: io.BytesIO) -> bytes:
    """
    Generates a thumbnail image for the first page of a PDF.
    """
    try:
        # Reset file pointer
        file.seek(0)
        # Convert first page of PDF to image
        images = convert_from_bytes(file.read(), first_page=1, last_page=1, size=(200, 200))
        if images:
            img_buffer = io.BytesIO()
            images[0].save(img_buffer, format='PNG')
            img_bytes = img_buffer.getvalue()
            logger.info("Generated PDF thumbnail.")
            return img_bytes
        else:
            logger.warning("No images generated for PDF thumbnail.")
            return None
    except Exception as e:
        st.warning(f"Failed to generate thumbnail for PDF: {e}")
        logger.error(f"Failed to generate thumbnail for PDF: {e}")
        return None

def get_thumbnail(file_type: str, file: io.BytesIO) -> bytes:
    """
    Retrieves a thumbnail based on the file type.
    """
    if file_type == "pdf":
        return generate_pdf_thumbnail(file)
    elif file_type == "pptx":
        # Placeholder for PPTX thumbnail generation
        return None
    else:
        return None
